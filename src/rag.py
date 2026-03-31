# ─────────────────────────────────────────────────
# TROY — RAG (Retrieval Augmented Generation)
# Infima Foundation A.C.
# Soporta: PDF, DOCX, TXT, XLSX, PPTX, JPG, PNG
# ─────────────────────────────────────────────────

import os
import sys
sys.path.insert(0, os.path.dirname(__file__))

from pypdf import PdfReader
from docx import Document as DocxDocument
import chromadb
from chromadb.utils import embedding_functions

CARPETA_DOCS = os.path.join(os.path.dirname(__file__), "..", "infima")
DB_PATH      = os.path.join(os.path.dirname(__file__), "..", "troy_rag_db")
COLECCION    = "infima_docs"

EXTENSIONES = {".pdf", ".docx", ".txt", ".xlsx", ".xls",
               ".pptx", ".ppt", ".csv", ".md", ".jpg",
               ".jpeg", ".png", ".webp"}

embedding_fn = embedding_functions.SentenceTransformerEmbeddingFunction(
    model_name="all-MiniLM-L6-v2"
)

cliente_chroma = chromadb.PersistentClient(path=DB_PATH)
coleccion = cliente_chroma.get_or_create_collection(
    name=COLECCION,
    embedding_function=embedding_fn
)

# ── LECTORES POR FORMATO ─────────────────────────

def leer_pdf(ruta: str) -> str:
    try:
        reader = PdfReader(ruta)
        return "".join(p.extract_text() or "" for p in reader.pages).strip()
    except Exception as e:
        print(f"Error PDF {ruta}: {e}")
        return ""

def leer_docx(ruta: str) -> str:
    try:
        doc = DocxDocument(ruta)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        print(f"Error DOCX {ruta}: {e}")
        return ""

def leer_txt(ruta: str) -> str:
    try:
        with open(ruta, "r", encoding="utf-8", errors="ignore") as f:
            return f.read().strip()
    except Exception as e:
        print(f"Error TXT {ruta}: {e}")
        return ""

def leer_csv(ruta: str) -> str:
    try:
        import csv
        filas = []
        with open(ruta, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for fila in reader:
                filas.append(", ".join(fila))
        return "\n".join(filas)
    except Exception as e:
        print(f"Error CSV {ruta}: {e}")
        return ""

def leer_xlsx(ruta: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(ruta, data_only=True)
        texto = []
        for hoja in wb.worksheets:
            texto.append(f"[Hoja: {hoja.title}]")
            for fila in hoja.iter_rows(values_only=True):
                celda = " | ".join(str(c) for c in fila if c is not None)
                if celda.strip():
                    texto.append(celda)
        return "\n".join(texto)
    except Exception as e:
        print(f"Error XLSX {ruta}: {e}")
        return ""

def leer_pptx(ruta: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(ruta)
        texto = []
        for i, slide in enumerate(prs.slides):
            texto.append(f"[Slide {i+1}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texto.append(shape.text.strip())
        return "\n".join(texto)
    except Exception as e:
        print(f"Error PPTX {ruta}: {e}")
        return ""

def leer_imagen(ruta: str) -> str:
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(ruta)
        texto = pytesseract.image_to_string(img, lang="spa+eng")
        return texto.strip()
    except Exception as e:
        print(f"Error imagen {ruta}: {e}")
        return ""

def leer_documento(ruta: str) -> str:
    ext = os.path.splitext(ruta)[1].lower()
    if ext == ".pdf":
        return leer_pdf(ruta)
    elif ext == ".docx":
        return leer_docx(ruta)
    elif ext in {".txt", ".md"}:
        return leer_txt(ruta)
    elif ext == ".csv":
        return leer_csv(ruta)
    elif ext in {".xlsx", ".xls"}:
        return leer_xlsx(ruta)
    elif ext in {".pptx", ".ppt"}:
        return leer_pptx(ruta)
    elif ext in {".jpg", ".jpeg", ".png", ".webp"}:
        return leer_imagen(ruta)
    return ""

# ── INDEXACIÓN ───────────────────────────────────

def dividir_en_chunks(texto: str, tamanio: int = 1000, overlap: int = 150) -> list:
    chunks = []
    inicio = 0
    while inicio < len(texto):
        fin = inicio + tamanio
        chunk = texto[inicio:fin]
        if chunk.strip():
            chunks.append(chunk)
        inicio = fin - overlap
    return chunks

def indexar_documentos():
    if not os.path.exists(CARPETA_DOCS):
        print(f"Carpeta no encontrada: {CARPETA_DOCS}")
        return

    archivos = [
        f for f in os.listdir(CARPETA_DOCS)
        if os.path.splitext(f)[1].lower() in EXTENSIONES
    ]

    if not archivos:
        print("No hay documentos para indexar.")
        return

    docs_existentes = set(coleccion.get()["ids"]) if coleccion.count() > 0 else set()
    nuevos = 0

    for archivo in archivos:
        ruta = os.path.join(CARPETA_DOCS, archivo)
        texto = leer_documento(ruta)

        if not texto:
            continue

        chunks = dividir_en_chunks(texto)

        for i, chunk in enumerate(chunks):
            doc_id = f"{archivo}_chunk_{i}"
            if doc_id in docs_existentes:
                continue
            coleccion.add(
                documents=[chunk],
                ids=[doc_id],
                metadatas=[{"archivo": archivo, "chunk": i}]
            )
            nuevos += 1

    print(f"RAG: {nuevos} fragmentos indexados de {len(archivos)} documentos.")

# ── BÚSQUEDA ─────────────────────────────────────

def buscar_contexto(pregunta: str, n_resultados: int = 5) -> str:
    if coleccion.count() == 0:
        return ""

    resultados = coleccion.query(
        query_texts=[pregunta],
        n_results=min(n_resultados, coleccion.count())
    )

    if not resultados["documents"] or not resultados["documents"][0]:
        return ""

    contexto = ""
    for i, doc in enumerate(resultados["documents"][0]):
        archivo = resultados["metadatas"][0][i]["archivo"]
        contexto += f"\n[Fragmento de {archivo}]:\n{doc}\n"

    return contexto.strip()

# Indexar al importar
indexar_documentos()